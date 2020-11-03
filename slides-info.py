from pptx import Presentation
import re

prs = Presentation('PMC - Quantitative Risk Update_20190115 - Copy.pptx')
notes_master = prs.notes_master

text_runs = []

for x, slide in enumerate(prs.slides):
    y = slide.notes_slide.notes_text_frame.text
    z = re.findall('(Est time:)(.*)',y)
    ##print(x,'-----------------------------------------------------\n',y,'\n',z,'\n')
    print(x+1,z)
